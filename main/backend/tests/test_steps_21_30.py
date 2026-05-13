from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import patch

from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from app.api.app import create_app
from app.config import settings


def _build_docx_bytes(text: str) -> bytes:
    with TemporaryDirectory() as tmp:
        p = Path(tmp) / "sample.docx"
        doc = Document()
        doc.add_paragraph(text)
        doc.save(str(p))
        return p.read_bytes()


def _build_pptx_bytes(title: str, bullet: str) -> bytes:
    with TemporaryDirectory() as tmp:
        p = Path(tmp) / "sample.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = bullet
        prs.save(str(p))
        return p.read_bytes()


def _build_pdf_bytes(text: str) -> bytes:
    with TemporaryDirectory() as tmp:
        p = Path(tmp) / "sample.pdf"
        c = canvas.Canvas(str(p), pagesize=letter)
        c.drawString(100, 750, text)
        c.save()
        return p.read_bytes()


def _register_and_login(client: TestClient):
    client.post("/v1/auth/register", json={"username": "tester", "password": "123456"})
    login = client.post("/v1/auth/login", json={"account": "tester", "password": "123456"})
    assert login.status_code == 200
    token = login.json()["data"]["token"]
    user_id = login.json()["data"]["user_id"]
    return token, user_id


def test_step21_to_step30_core_flows():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            token, user_id = _register_and_login(client)
            auth = {"Authorization": f"Bearer {token}"}

            me = client.get("/v1/auth/me", headers=auth)
            assert me.status_code == 200
            upd = client.put("/v1/auth/profile", json={"username": "tester2"}, headers=auth)
            assert upd.status_code == 200
            pwd = client.put(
                "/v1/auth/password",
                json={"old_password": "123456", "new_password": "abcdef"},
                headers=auth,
            )
            assert pwd.status_code == 200

            provider = client.post(
                "/v1/providers",
                json={
                    "user_id": user_id,
                    "provider_type": "ollama",
                    "display_name": "Local Ollama",
                    "base_url": "http://localhost:11434",
                    "model_name": "qwen2.5:7b",
                    "is_default": True,
                },
                headers=auth,
            )
            assert provider.status_code == 200
            provider_list = client.get(f"/v1/providers/{user_id}", headers=auth)
            assert provider_list.status_code == 200
            assert len(provider_list.json()["data"]["items"]) >= 1
            provider_test = client.post(
                "/v1/providers/test",
                json={"provider_type": "openai_compatible", "base_url": "https://api.deepseek.com", "model_name": "deepseek-chat"},
                headers=auth,
            )
            assert provider_test.status_code == 200

            skills_all = client.get("/v1/skills")
            assert skills_all.status_code == 200
            skills_resolved = client.get("/v1/skills/resolve/ppt/generation")
            assert skills_resolved.status_code == 200

            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Generate 10 slides",
                    "pages": 10,
                    "style": "academic_simple",
                    "language": "zh-CN",
                },
                headers=auth,
            )
            task_id = create.json()["data"]["task_id"]

            with client.websocket_connect(f"/ws/tasks/{task_id}?token={token}") as ws:
                payload = ws.receive_json()
                assert payload["task_id"] == task_id

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(b"hello task flow"), "text/plain")},
                headers=auth,
            )
            assert upload.status_code == 200
            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False}, headers=auth)
            assert parse.status_code == 200
            run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False}, headers=auth)
            assert run.status_code == 200

            versions = client.get(f"/v1/tasks/{task_id}/versions", headers=auth)
            assert versions.status_code == 200
            assert len(versions.json()["data"]["items"]) == 1

            for i in range(4):
                llm_revision_json = (
                    "{"
                    f"\"title\":\"Revised Slide {i}\","
                    "\"bullets\":[\"Updated point A\",\"Updated point B\",\"Updated point C\"],"
                    "\"notes\":\"Revised notes for this slide.\","
                    "\"image_placeholders\":[{\"label\":\"Revision figure\",\"source\":\"web_search\"}]"
                    "}"
                )
                with patch(
                    "app.services.llm_runtime.text_generator.LLMTextGenerator.generate",
                    return_value=llm_revision_json,
                ) as mocked_llm:
                    rev = client.post(
                        f"/v1/tasks/{task_id}/revisions",
                        json={"page_index": 2, "instruction": f"rev-{i}"},
                        headers=auth,
                    )
                    assert rev.status_code == 200
                    assert mocked_llm.called

            versions2 = client.get(f"/v1/tasks/{task_id}/versions", headers=auth)
            items = versions2.json()["data"]["items"]
            assert len(items) == 5

            compare = client.get(f"/v1/tasks/{task_id}/versions/compare?from_version=1&to_version=5", headers=auth)
            assert compare.status_code == 200
            assert compare.json()["data"]["changed_page_count"] >= 1

            rollback = client.post(f"/v1/tasks/{task_id}/versions/rollback/2", json={}, headers=auth)
            assert rollback.status_code == 200
            assert rollback.json()["data"]["new_version"] == 6

            dl_v2 = client.get(f"/v1/tasks/{task_id}/download/2", headers=auth)
            assert dl_v2.status_code == 200
            assert dl_v2.json()["data"]["exists"] is True


def test_step22_min_regression_set():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            token, user_id = _register_and_login(client)
            auth = {"Authorization": f"Bearer {token}"}

            valid_samples = [
                ("ok1.pdf", _build_pdf_bytes("pdf sample 1")),
                ("ok2.pdf", _build_pdf_bytes("pdf sample 2")),
                ("ok3.docx", _build_docx_bytes("docx sample 3")),
                ("ok4.docx", _build_docx_bytes("docx sample 4")),
                ("ok5.doc", b"legacy doc plain bytes"),
                ("ok6.txt", b"plain txt"),
                ("ok7.pptx", _build_pptx_bytes("title", "bullet")),
                ("ok8.ppt", _build_pptx_bytes("title2", "bullet2")),
            ]
            invalid_samples = [
                ("bad1.exe", b"binary"),
                ("bad2.txt", b""),  # empty
                ("bad3.txt", b"x" * (2 * 1024 * 1024)),  # oversize under lowered limit
            ]
            damaged_sample = ("bad4.docx", b"not-a-real-docx")

            original_limit = settings.max_upload_size_bytes
            settings.max_upload_size_bytes = 1024 * 1024
            try:
                for name, content in valid_samples:
                    create = client.post(
                        "/v1/tasks",
                        json={
                            "user_id": user_id,
                            "user_requirement": "regression valid",
                            "pages": 10,
                            "style": "academic_simple",
                            "language": "zh-CN",
                        },
                        headers=auth,
                    )
                    task_id = create.json()["data"]["task_id"]
                    upload = client.post(
                        f"/v1/tasks/{task_id}/upload",
                        files={"upload": (name, BytesIO(content), "application/octet-stream")},
                        headers=auth,
                    )
                    assert upload.status_code == 200

                for name, content in invalid_samples:
                    create = client.post(
                        "/v1/tasks",
                        json={
                            "user_id": user_id,
                            "user_requirement": "regression invalid",
                            "pages": 10,
                            "style": "academic_simple",
                            "language": "zh-CN",
                        },
                        headers=auth,
                    )
                    task_id = create.json()["data"]["task_id"]
                    upload = client.post(
                        f"/v1/tasks/{task_id}/upload",
                        files={"upload": (name, BytesIO(content), "application/octet-stream")},
                        headers=auth,
                    )
                    assert upload.status_code == 400

                create = client.post(
                    "/v1/tasks",
                    json={
                        "user_id": user_id,
                        "user_requirement": "regression damaged",
                        "pages": 10,
                        "style": "academic_simple",
                        "language": "zh-CN",
                    },
                    headers=auth,
                )
                task_id = create.json()["data"]["task_id"]
                upload = client.post(
                    f"/v1/tasks/{task_id}/upload",
                    files={"upload": (damaged_sample[0], BytesIO(damaged_sample[1]), "application/octet-stream")},
                    headers=auth,
                )
                assert upload.status_code == 200
                parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False}, headers=auth)
                assert parse.status_code == 400
            finally:
                settings.max_upload_size_bytes = original_limit


def test_revision_without_page_index_can_auto_select_multiple_slides():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            token, user_id = _register_and_login(client)
            auth = {"Authorization": f"Bearer {token}"}
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Generate 8 slides about AI governance",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
                headers=auth,
            )
            task_id = create.json()["data"]["task_id"]
            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(b"AI governance baseline content"), "text/plain")},
                headers=auth,
            )
            assert upload.status_code == 200
            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False}, headers=auth)
            assert parse.status_code == 200
            run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False}, headers=auth)
            assert run.status_code == 200

            selector_json = "{\"page_indexes\":[2,3],\"reason\":\"scope spans two slides\"}"
            revised_2 = "{\"title\":\"Revised 2\",\"bullets\":[\"A\",\"B\",\"C\"],\"notes\":\"N2\",\"image_placeholders\":[]}"
            revised_3 = "{\"title\":\"Revised 3\",\"bullets\":[\"D\",\"E\",\"F\"],\"notes\":\"N3\",\"image_placeholders\":[]}"
            with patch(
                "app.services.llm_runtime.text_generator.LLMTextGenerator.generate",
                side_effect=[selector_json, revised_2, revised_3],
            ):
                rev = client.post(
                    f"/v1/tasks/{task_id}/revisions",
                    json={"instruction": "Please merge duplicated points and refine logic"},
                    headers=auth,
                )
                assert rev.status_code == 200
                revised_pages = rev.json()["data"].get("revised_pages", [])
                assert revised_pages == [2, 3]
