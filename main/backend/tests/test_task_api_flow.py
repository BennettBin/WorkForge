from io import BytesIO
import json
import shutil
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import patch

from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation

from app.api.app import create_app
from app.config import settings
from app.models.entities import LLMProviderConfig
from app.services.skill_runtime.executor import SkillExecutor


def _register_and_login(client: TestClient, username: str = "task_flow_user") -> tuple[str, dict[str, str]]:
    client.post("/v1/auth/register", json={"username": username, "password": "123456"})
    login = client.post("/v1/auth/login", json={"account": username, "password": "123456"})
    assert login.status_code == 200
    data = login.json()["data"]
    auth = {"Authorization": f"Bearer {data['token']}"}
    client.headers.update(auth)
    return data["user_id"], auth


def test_task_api_create_upload_parse_run_flow():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_create_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Create 10-page slides from source file",
                    "template_choice": "system_default",
                    "pages": 10,
                    "style": "academic_simple",
                    "language": "zh-CN",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(b"hello workforge"), "text/plain")},
            )
            assert upload.status_code == 200
            assert upload.json()["data"]["file_type"] == "txt"

            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
            assert parse.status_code == 200
            assert parse.json()["data"]["parse_status"] == "success"
            vector_index_path = Path(temp_dir) / "vectors" / task_id / "index.json"
            assert vector_index_path.exists()
            payload = json.loads(vector_index_path.read_text(encoding="utf-8"))
            assert payload["chunk_count"] > 0

            run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
            assert run.status_code == 200
            assert run.json()["data"]["status"] == "completed"
            pptx_path = Path(run.json()["data"]["output_path"])
            assert pptx_path.exists()
            prs = Presentation(str(pptx_path))
            assert len(prs.slides) == 10

            get_task = client.get(f"/v1/tasks/{task_id}")
            assert get_task.status_code == 200
            assert get_task.json()["data"]["task"]["status"] == "completed"
            skill_names = [row["skill_name"] for row in get_task.json()["data"]["skill_calls"]]
            assert "find_skill" in skill_names
            events = get_task.json()["data"]["events"]
            assert any("template_script_render_succeeded" in e["message"] for e in events)

            download = client.get(f"/v1/tasks/{task_id}/download/latest")
            assert download.status_code == 200
            assert download.json()["data"]["exists"] is True


def test_task_api_reject_empty_upload():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_empty_upload_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Create slides",
                    "template_choice": "system_default",
                    "pages": 10,
                    "style": "academic_simple",
                    "language": "zh-CN",
                },
            )
            task_id = create.json()["data"]["task_id"]

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("empty.txt", BytesIO(b""), "text/plain")},
            )
            assert upload.status_code == 400
            assert upload.json()["error"]["code"] == "BAD_REQUEST"


def test_task_history_delete_removes_records_and_files():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_delete_history_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Create deletable history slides",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("delete-me.txt", BytesIO(b"delete history source"), "text/plain")},
            )
            assert upload.status_code == 200
            source_path = Path(upload.json()["data"]["file_path"])
            assert source_path.exists()

            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
            assert parse.status_code == 200
            parsed_path = Path(parse.json()["data"]["parsed_text_path"])
            assert parsed_path.exists()

            run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
            assert run.status_code == 200
            output_path = Path(run.json()["data"]["output_path"])
            assert output_path.exists()

            versions = client.get(f"/v1/tasks/{task_id}/versions")
            assert versions.status_code == 200
            assert len(versions.json()["data"]["items"]) >= 1

            delete = client.delete(f"/v1/tasks/{task_id}")
            assert delete.status_code == 200
            deleted = delete.json()["data"]["deleted"]
            assert deleted["task"] is True
            assert deleted["files"] == 1
            assert deleted["outputs"] >= 1
            assert not source_path.exists()
            assert not parsed_path.exists()
            assert not output_path.exists()

            listed = client.get(f"/v1/tasks/user/{user_id}")
            assert listed.status_code == 200
            assert all(row["task_id"] != task_id for row in listed.json()["data"]["items"])

            get_deleted = client.get(f"/v1/tasks/{task_id}")
            assert get_deleted.status_code == 400


def test_task_api_triggers_knowledge_search_skill_when_requested():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_search_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "请搜索最新研究并补充到PPT",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "zh-CN",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            original_execute = SkillExecutor.execute

            def _mock_execute(self, skill_name, payload):
                if skill_name == "knowledge_search":
                    return {"items": [{"title": "t", "url": "https://example.com/a", "snippet": "s", "content": "external web content for vector cache"}]}
                return original_execute(self, skill_name, payload)

            with patch("app.services.skill_runtime.executor.SkillExecutor.execute", new=_mock_execute):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
                assert run.status_code == 200

            detail = client.get(f"/v1/tasks/{task_id}")
            assert detail.status_code == 200
            skill_calls = detail.json()["data"]["skill_calls"]
            assert any(call["skill_name"] == "knowledge_search" for call in skill_calls)
            search_outputs = [
                json.loads(call["output"]) if isinstance(call.get("output"), str) else call.get("output", {})
                for call in skill_calls
                if call["skill_name"] == "knowledge_search"
            ]
            assert any(output.get("count", 0) > 0 for output in search_outputs)


def test_task_api_can_clear_vector_cache_manually():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_cache_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Create slides",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "zh-CN",
                },
            )
            task_id = create.json()["data"]["task_id"]

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(b"cache clear test input"), "text/plain")},
            )
            assert upload.status_code == 200
            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
            assert parse.status_code == 200

            vector_index_path = Path(temp_dir) / "vectors" / task_id / "index.json"
            assert vector_index_path.exists()

            clear = client.post(f"/v1/tasks/{task_id}/cache/clear", json={})
            assert clear.status_code == 200
            assert clear.json()["data"]["removed"] is True
            assert not vector_index_path.exists()


def test_task_api_no_source_file_can_still_generate_by_search():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_no_source_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Create PPT about AI trends without source file",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            with patch(
                "app.skills.knowledge_search.runtime.KnowledgeSearchService.search_and_extract",
                return_value=[{"title": "AI trend", "url": "https://example.com/trend", "snippet": "trend snippet", "content": "trend content"}],
            ):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
                assert run.status_code == 200
                assert run.json()["data"]["status"] == "completed"
                pptx_path = Path(run.json()["data"]["output_path"])
                assert pptx_path.exists()

            detail = client.get(f"/v1/tasks/{task_id}")
            assert detail.status_code == 200
            events = detail.json()["data"]["events"]
            assert any("no_source_file_detected" in e["message"] for e in events)


def test_task_api_no_source_file_search_empty_still_generates_with_fallback():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_no_source_empty_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Create PPT about AI governance without source file",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            with patch(
                "app.skills.knowledge_search.runtime.KnowledgeSearchService.search_and_extract",
                return_value=[],
            ):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
                assert run.status_code == 200
                assert run.json()["data"]["status"] == "completed"
                assert Path(run.json()["data"]["output_path"]).exists()

            detail = client.get(f"/v1/tasks/{task_id}")
            assert detail.status_code == 200
            events = detail.json()["data"]["events"]
            assert any("no_source_file_forced_search_empty;fallback_to_requirement_context" in e["message"] for e in events)


def test_task_api_exports_image_placeholder_metadata():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_image_slots_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Make a research presentation with figure slots",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            task_id = create.json()["data"]["task_id"]
            file_bytes = (
                b"Figure 1: pipeline architecture from source_file\n"
                b"Methods and results\n"
                b"Figure 2: ablation chart from source_file\n"
            )
            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(file_bytes), "text/plain")},
            )
            assert upload.status_code == 200
            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
            assert parse.status_code == 200

            run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
            assert run.status_code == 200
            slides_path = Path(run.json()["data"]["slides_path"])
            slides = json.loads(slides_path.read_text(encoding="utf-8"))
            assert any(len(s.get("image_placeholders", [])) > 0 for s in slides if s.get("kind") == "content")


def test_task_parse_uses_ollama_embedding_config_when_default_provider_is_ollama():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_ollama_embedding_user")
            # seed default provider for same user
            app.state.repositories.providers.upsert(
                LLMProviderConfig(
                    provider_id="p-ollama",
                    user_id=user_id,
                    provider_type="ollama",
                    display_name="Ollama",
                    base_url="http://localhost:11434",
                    model_name="qwen3:8b",
                    embedding_model="qwen3-embedding:8b",
                    is_default=True,
                )
            )

            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "embedding config test",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]
            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(b"embedding test content"), "text/plain")},
            )
            assert upload.status_code == 200

            with patch(
                "app.services.vector_store.index_service.VectorIndexService._vectorize_ollama",
                return_value={"0": 1.0},
            ) as mocked:
                parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
                assert parse.status_code == 200
                assert mocked.called

            vector_index_path = Path(temp_dir) / "vectors" / task_id / "index.json"
            payload = json.loads(vector_index_path.read_text(encoding="utf-8"))
            assert payload.get("vectorizer", {}).get("type") == "ollama"
            assert payload.get("vectorizer", {}).get("model") == "qwen3-embedding:8b"


def test_task_api_llm_failure_falls_back_when_not_required():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_llm_fallback_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Create slides with llm",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            task_id = create.json()["data"]["task_id"]
            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(b"llm fallback source"), "text/plain")},
            )
            assert upload.status_code == 200
            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
            assert parse.status_code == 200

            with patch(
                "app.services.llm_runtime.text_generator.LLMTextGenerator.generate",
                side_effect=RuntimeError("forced llm failure for test"),
            ):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False, "require_llm": False})
                assert run.status_code == 200
                assert run.json()["data"]["status"] == "completed"
                llm_debug = run.json()["data"]["llm_debug"]
                assert llm_debug["attempted"] is True
                assert llm_debug["succeeded"] is False
                assert llm_debug["failed_reason"] is not None


def test_task_api_llm_failure_blocks_when_required():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_llm_required_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "user_requirement": "Create slides with llm required",
                    "template_choice": "system_default",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            task_id = create.json()["data"]["task_id"]
            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(b"llm required source"), "text/plain")},
            )
            assert upload.status_code == 200
            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
            assert parse.status_code == 200

            with patch(
                "app.services.llm_runtime.text_generator.LLMTextGenerator.generate",
                side_effect=RuntimeError("forced llm failure for test"),
            ):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False, "require_llm": True})
                assert run.status_code == 400
                body = run.json()
                assert body["error"]["code"] == "BAD_REQUEST"
                assert "LLM required" in body["error"]["message"]


def test_task_api_extended_task_types_generate_markdown_outputs():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_extended_types_user")
            for task_type in ["report", "wechat_post", "data_analysis", "code_doc", "paper_assistant"]:
                create = client.post(
                    "/v1/tasks",
                    json={
                        "user_id": user_id,
                        "task_type": task_type,
                        "user_requirement": f"Generate {task_type} content",
                        "pages": 8,
                        "style": "academic_simple",
                        "language": "en-US",
                    },
                )
                assert create.status_code == 200
                task_id = create.json()["data"]["task_id"]

                with patch(
                    "app.services.llm_runtime.text_generator.LLMTextGenerator.generate",
                    return_value=(
                            (
                            "# report\n\n## Summary\nGenerated by llm with detailed context and expanded explanation for reviewers.\n\n## Findings\nPoint A with evidence.\nPoint B with trend.\n\n## Recommendations\nPoint C with action plan.\n"
                            if task_type == "report"
                            else "# wechat_post\n\n## Summary\nGenerated by llm with clear narrative and audience-oriented language.\n\n## Findings\nPoint A for readers.\nPoint B for context.\n\n## Recommendations\nPoint C as next step.\n"
                            if task_type == "wechat_post"
                            else "# data_analysis\n\n## Cleaning\nStep A normalization.\nStep B missing-value strategy.\n\n## Findings\nPoint B distribution.\nPoint C correlation.\n\n## Recommendations\nPoint D chart plan.\n"
                            if task_type == "data_analysis"
                            else "# code_doc\n\n## Quick Start\nStep A install dependencies.\nStep B run service.\n\n## API\nEndpoint B request schema.\nEndpoint C response example.\n\n## Recommendations\nPoint D maintenance.\n"
                            if task_type == "code_doc"
                            else "# paper_assistant\n\n## Abstract\nDraft A with motivation and method overview.\n\n## Revision\nSuggestion B for clarity.\nSuggestion C for evidence linkage.\n\n## Findings\nPoint D contribution scope.\n"
                        )
                    ),
                ), patch(
                    "app.skills.knowledge_search.runtime.KnowledgeSearchService.search_and_extract",
                    return_value=[{"title": "t", "url": "https://example.com", "snippet": "s", "content": "c"}],
                ):
                    run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
                    assert run.status_code == 200
                    assert run.json()["data"]["status"] == "completed"
                    assert run.json()["data"]["task_type"] == task_type
                    output_path = Path(run.json()["data"]["output_path"])
                    assert output_path.exists()
                    assert output_path.suffix == ".md"
                detail = client.get(f"/v1/tasks/{task_id}")
                assert detail.status_code == 200
                skill_names = [row["skill_name"] for row in detail.json()["data"]["skill_calls"]]
                assert "find_skill" in skill_names
                assert "skill_registry_resolve" in skill_names
                assert any(name in skill_names for name in {
                    "report_generation",
                    "wechat_post_generation",
                    "data_analysis",
                    "code_doc_generation",
                    "paper_assistant_generation",
                })


def test_task_api_continues_when_find_skill_returns_no_match():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_find_skill_no_match_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "report",
                    "user_requirement": "Generate weekly report",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]
            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.txt", BytesIO(b"weekly report source content"), "text/plain")},
            )
            assert upload.status_code == 200
            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
            assert parse.status_code == 200

            original_execute = SkillExecutor.execute

            def _mock_execute(self, skill_name, payload):
                if skill_name == "find_skill":
                    return {"task_type": "report", "stage": "generation", "has_match": False, "matched_skills": [], "all_candidates": []}
                return original_execute(self, skill_name, payload)

            with patch("app.services.skill_runtime.executor.SkillExecutor.execute", new=_mock_execute):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
                assert run.status_code == 200
                assert run.json()["data"]["status"] == "completed"
                assert Path(run.json()["data"]["output_path"]).exists()


def test_non_ppt_revision_uses_markdown_pipeline_instead_of_slide_json():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_non_ppt_revision_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "wechat_post",
                    "user_requirement": "Write a WeChat post about AI agents",
                    "pages": 10,
                    "style": "academic_simple",
                    "language": "zh-CN",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            with patch(
                "app.services.llm_runtime.text_generator.LLMTextGenerator.generate",
                side_effect=[
                    "# Wechat Post\n\n## Summary\nThis post explains practical value, risks, and adoption steps.\n\n## Findings\nPoint A: real scenarios.\nPoint B: implementation checklist.\nPoint C: common pitfalls.\n\n## Recommendations\nFollow and share your experience.",
                    "# Wechat Post\n\n## Summary\nThis revised version is more concise and emphasizes action.\n\n## Findings\nPoint A: practical scenario in one sentence.\nPoint B: checklist with clear priority.\nPoint C: concise risk reminder.\n\n## Recommendations\nComment your use case and follow for templates.",
                ],
            ):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
                assert run.status_code == 200
                assert run.json()["data"]["status"] == "completed"
                rev = client.post(
                    f"/v1/tasks/{task_id}/revisions",
                    json={"instruction": "Make it more concise and stronger CTA"},
                )
                assert rev.status_code == 200
                assert rev.json()["data"]["new_version"] == 2

            versions = client.get(f"/v1/tasks/{task_id}/versions")
            assert versions.status_code == 200
            items = versions.json()["data"]["items"]
            assert len(items) == 2
            latest_path = Path(items[-1]["file_path"])
            assert latest_path.exists()
            assert latest_path.suffix == ".md"


def test_data_analysis_task_accepts_xlsx_and_exports_docx_report():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_data_analysis_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "data_analysis",
                    "user_requirement": "Analyze cate distribution and return a Word report with chart.",
                    "pages": 10,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            wb = Workbook()
            ws = wb.active
            ws.append(["company", "cate", "revenue"])
            ws.append(["A", "Manufacturing", 100])
            ws.append(["B", "Manufacturing", 120])
            ws.append(["C", "Service", 80])
            ws.append(["D", "Tech", 140])
            stream = BytesIO()
            wb.save(stream)
            stream.seek(0)

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("sample.xlsx", stream, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
            )
            assert upload.status_code == 200

            parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False})
            assert parse.status_code == 200

            with patch(
                "app.services.llm_runtime.text_generator.LLMTextGenerator.generate",
                return_value=(
                    "# Data Analysis\n\n## Cleaning\nMissing values checked.\nType normalized.\n\n"
                    "## Findings\nManufacturing dominates cate.\nService and Tech are smaller.\n\n"
                    "## Recommendations\nUse balanced sampling in next-stage modeling.\n"
                ),
            ):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
                assert run.status_code == 200
                out = Path(run.json()["data"]["output_path"])
                assert out.exists()
                assert out.suffix == ".docx"


def test_ppt_template_extraction_and_template_listing():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        templates_root = Path(__file__).resolve().parents[1] / "app" / "templates" / "ppt"
        template_dir = templates_root / "OncologyTemplate"
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_template_extract_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "ppt",
                    "user_requirement": "请提取模板并保存模板\nTemplate_Name=OncologyTemplate",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "zh-CN",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            if slide.shapes.title is not None:
                slide.shapes.title.text = "Template Title"
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "Template Subtitle"
            stream = BytesIO()
            prs.save(stream)
            stream.seek(0)

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("template_source.pptx", stream, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            )
            assert upload.status_code == 200

            run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
            assert run.status_code == 200
            data = run.json()["data"]
            assert data["template_extracted"] is True
            assert data["template_name"] == "OncologyTemplate"
            assert Path(data["template_file"]).exists()
            assert Path(data["metadata_file"]).exists()

            templates = client.get("/v1/tasks/ppt/templates")
            assert templates.status_code == 200
            items = templates.json()["data"]["items"]
            assert any(item["name"] == "OncologyTemplate" for item in items)
        if template_dir.exists():
            shutil.rmtree(template_dir)


def test_ppt_template_generation_recovery_and_resume_flow():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        templates_root = Path(__file__).resolve().parents[1] / "app" / "templates" / "ppt"
        template_dir = templates_root / "RecoveryTemplate"
        with TestClient(app) as client:
            user_id, auth = _register_and_login(client, "ppt_recovery_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "ppt",
                    "user_requirement": "extract template for ppt\nTemplate_Name=RecoveryTemplate\nForceInvalidBundle=true",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
                headers=auth,
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[0])
            stream = BytesIO()
            prs.save(stream)
            stream.seek(0)

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("template_source.pptx", stream, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                headers=auth,
            )
            assert upload.status_code == 200

            run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False}, headers=auth)
            assert run.status_code == 200
            data = run.json()["data"]
            assert data["status"] == "requires_user_completion"
            assert data.get("resume_token")
            assert data["missing_items"] == ["template.rules.json"]
            assert data["suggested_values"].get("rules.rules")
            assert data["bundle_validation"]["ok"] is False

            recovery = client.get(f"/v1/tasks/{task_id}/template-generation/recovery", headers=auth)
            assert recovery.status_code == 200
            rec = recovery.json()["data"]
            assert rec["has_recovery"] is True
            token = rec["resume_token"]
            assert rec["missing_items"] == ["template.rules.json"]
            detail = client.get(f"/v1/tasks/{task_id}", headers=auth)
            assert detail.status_code == 200
            detail_data = detail.json()["data"]
            assert any(e["stage"] == "requires_user_completion" for e in detail_data["events"])
            assert any(c["skill_name"] == "template_bundle_validation" for c in detail_data["skill_calls"])

            resume = client.post(
                f"/v1/tasks/{task_id}/template-generation/resume",
                json={
                    "resume_token": token,
                    "user_filled_fields": {
                        "text_style.title.size_pt": "32",
                        "text_style.body.size_pt": "18",
                        "rules.schema_version": "v1",
                        "rules.rules": "[]",
                    },
                },
                headers=auth,
            )
            assert resume.status_code == 200
            assert resume.json()["data"]["status"] == "completed"
            assert resume.json()["data"]["applied_rules"] == 0
            after_resume = client.get(f"/v1/tasks/{task_id}", headers=auth)
            after_data = after_resume.json()["data"]
            assert any("template_generation_resume_started" in e["message"] for e in after_data["events"])
            assert any("template_generation_resume_completed" in e["message"] for e in after_data["events"])
            assert len(after_data["outputs"]) == 1
        if template_dir.exists():
            shutil.rmtree(template_dir)


def test_ppt_template_generation_aux_llm_uses_user_default_model():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        templates_root = Path(__file__).resolve().parents[1] / "app" / "templates" / "ppt"
        template_dir = templates_root / "UserModelTemplate"
        calls: list[dict] = []

        def fake_generate(self, **kwargs):
            calls.append(kwargs)
            prompt = str(kwargs.get("prompt", ""))
            if "Choose necessary skills" in prompt:
                return json.dumps({"skills": ["ppt_template_generation"]})
            return json.dumps({"rules.schema_version": "v1", "rules.rules": "[]"})

        with TestClient(app) as client:
            user_id, auth = _register_and_login(client, "ppt_template_model_user")
            app.state.repositories.providers.upsert(
                LLMProviderConfig(
                    provider_id="user-provider",
                    user_id=user_id,
                    provider_type="openai_compatible",
                    display_name="User Model",
                    base_url="https://user-model.example/v1",
                    api_key_encrypted="user-secret",
                    model_name="user-configured-model",
                    chat_model="user-configured-model",
                    timeout_seconds=77,
                    is_default=True,
                )
            )
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "ppt",
                    "user_requirement": "extract template for ppt\nTemplate_Name=UserModelTemplate\nForceInvalidBundle=true",
                },
                headers=auth,
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[0])
            stream = BytesIO()
            prs.save(stream)
            stream.seek(0)
            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("template_source.pptx", stream, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                headers=auth,
            )
            assert upload.status_code == 200

            with patch("app.services.llm_runtime.text_generator.LLMTextGenerator.generate", fake_generate):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False}, headers=auth)

            assert run.status_code == 200
            data = run.json()["data"]
            assert data["status"] == "requires_user_completion"
            assert data["suggested_values"]["rules.rules"] == "[]"
            assert calls
            assert all(call["provider_type"] == "openai_compatible" for call in calls)
            assert all(call["base_url"] == "https://user-model.example/v1" for call in calls)
            assert all(call["model_name"] == "user-configured-model" for call in calls)
            assert all(call["api_key"] == "user-secret" for call in calls)

            detail = client.get(f"/v1/tasks/{task_id}", headers=auth)
            assert detail.status_code == 200
            skill_calls = detail.json()["data"]["skill_calls"]
            assert any(c["skill_name"] == "ppt_template_generation_llm" for c in skill_calls)
        if template_dir.exists():
            shutil.rmtree(template_dir)


def test_list_ppt_templates_returns_only_valid_bundles():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        templates_root = Path(__file__).resolve().parents[1] / "app" / "templates" / "ppt"
        valid_dir = templates_root / "valid_bundle_for_test"
        invalid_dir = templates_root / "invalid_bundle_for_test"
        legacy_dir = templates_root / "legacy_recovery_bundle_for_test"
        try:
            valid_dir.mkdir(parents=True, exist_ok=True)
            invalid_dir.mkdir(parents=True, exist_ok=True)
            legacy_dir.mkdir(parents=True, exist_ok=True)
            (valid_dir / "template.pptx").write_bytes(b"pptx")
            (valid_dir / "render_from_template.py").write_text(
                (
                    "from pathlib import Path\n"
                    "from pptx import Presentation\n\n"
                    "def render(payload, output_path, template_path, meta, rules):\n"
                    "    prs = Presentation(str(template_path))\n"
                    "    while len(prs.slides) > 0:\n"
                    "        slide_id = prs.slides._sldIdLst[0]\n"
                    "        prs.part.drop_rel(slide_id.rId)\n"
                    "        prs.slides._sldIdLst.remove(slide_id)\n"
                    "    for idx, slide_payload in enumerate(payload.get('slides', []), start=1):\n"
                    "        layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None\n"
                    "        slide = prs.slides.add_slide(layout)\n"
                    "        if slide.shapes.title is not None:\n"
                    "            slide.shapes.title.text = str(slide_payload.get('title', f'Slide {idx}'))\n"
                    "    out = Path(output_path)\n"
                    "    out.parent.mkdir(parents=True, exist_ok=True)\n"
                    "    prs.save(str(out))\n"
                    "    return str(out)\n"
                ),
                encoding="utf-8",
            )
            (valid_dir / "template.meta.json").write_text(
                json.dumps(
                    {
                        "schema_version": "v1",
                        "template_type": "ppt",
                        "template_name": "valid_bundle_for_test",
                        "slide_size": {"width_inches": 13.333, "height_inches": 7.5, "aspect_ratio": "16:9"},
                        "theme": {"name": "default", "palette": {"primary": "#000"}},
                        "layout_map": {"cover": "Title Slide", "content": "Title and Content", "summary": "Title and Content"},
                        "text_style": {"title": {"font": "Arial", "size_pt": 32}, "body": {"font": "Arial", "size_pt": 18}},
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            (valid_dir / "template.rules.json").write_text(json.dumps({"schema_version": "v1", "rules": []}), encoding="utf-8")

            # invalid bundle: missing rules and script
            (invalid_dir / "template.pptx").write_bytes(b"pptx")
            (invalid_dir / "template.meta.json").write_text(json.dumps({"schema_version": "v1", "template_type": "ppt"}), encoding="utf-8")

            for filename in ("template.pptx", "render_from_template.py", "template.meta.json", "template.rules.json"):
                source = valid_dir / filename
                target = legacy_dir / filename
                if source.suffix == ".pptx":
                    target.write_bytes(source.read_bytes())
                else:
                    target.write_text(source.read_text(encoding="utf-8"), encoding="utf-8")
            (legacy_dir / "template.recovery.json").write_text(json.dumps({"missing_items": []}), encoding="utf-8")

            with TestClient(app) as client:
                user_id, auth = _register_and_login(client, "list_templates_user")
                res = client.get("/v1/tasks/ppt/templates", headers=auth)
                assert res.status_code == 200
                items = res.json()["data"]["items"]
                names = [x["name"] for x in items]
                assert "system_default" in names
                assert "valid_bundle_for_test" in names
                assert "invalid_bundle_for_test" not in names
                assert "legacy_recovery_bundle_for_test" not in names
                row = next(x for x in items if x["name"] == "valid_bundle_for_test")
                assert row["is_valid"] is True
                assert row["missing_files"] == []
                assert row["forbidden_files"] == []
                assert row["schema_version"] == "v1"
                default_row = next(x for x in items if x["name"] == "system_default")
                assert default_row["is_valid"] is True
                assert default_row["missing_files"] == []
                assert default_row["forbidden_files"] == []
                assert default_row["schema_version"] == "v1"

                diagnostic = client.get("/v1/tasks/ppt/templates?include_invalid=true", headers=auth)
                assert diagnostic.status_code == 200
                diagnostic_items = diagnostic.json()["data"]["items"]
                invalid_row = next(x for x in diagnostic_items if x["name"] == "invalid_bundle_for_test")
                assert invalid_row["is_valid"] is False
                assert "template.rules.json" in invalid_row["missing_files"]
                assert invalid_row["validation_errors"]
                legacy_row = next(x for x in diagnostic_items if x["name"] == "legacy_recovery_bundle_for_test")
                assert legacy_row["is_valid"] is False
                assert legacy_row["missing_files"] == []
                assert legacy_row["forbidden_files"] == ["template.recovery.json"]
                assert "legacy/recovery artifact" in legacy_row["validation_errors"][0]
                create = client.post(
                    "/v1/tasks",
                    json={
                        "user_id": user_id,
                        "task_type": "ppt",
                        "user_requirement": "Build slides with legacy template.",
                        "template_choice": "legacy_recovery_bundle_for_test",
                        "pages": 8,
                        "style": "legacy_recovery_bundle_for_test",
                        "language": "en-US",
                    },
                    headers=auth,
                )
                assert create.status_code == 400
                assert "TemplateChoice points to invalid template bundle" in create.text
                assert "template.recovery.json" in create.text
        finally:
            if valid_dir.exists():
                shutil.rmtree(valid_dir)
            if invalid_dir.exists():
                shutil.rmtree(invalid_dir)
            if legacy_dir.exists():
                shutil.rmtree(legacy_dir)


def test_ppt_generation_fails_when_templatechoice_is_invalid():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, auth = _register_and_login(client, "invalid_templatechoice_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "ppt",
                    "user_requirement": "Build slides for demo.",
                    "template_choice": "non_existing_template",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
                headers=auth,
            )
            assert create.status_code == 400
            assert "TemplateChoice points to non-existent template" in create.text


def test_ppt_generation_defaults_to_system_templatechoice_when_missing():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, auth = _register_and_login(client, "missing_templatechoice_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "ppt",
                    "user_requirement": "Build slides for demo.",
                    "pages": 8,
                    "style": "system_default",
                    "language": "en-US",
                },
                headers=auth,
            )
            assert create.status_code == 200
            data = create.json()["data"]
            assert data["template_choice"] == "system_default"
            assert "TemplateChoice=system_default" in data["user_requirement"]


def test_ppt_generation_uses_full_path_template_choice_and_llm_content():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        template_dir = Path(__file__).resolve().parents[1] / "app" / "templates" / "ppt" / "ppt1"
        calls: list[str] = []

        def fake_generate(self, **kwargs):
            prompt = str(kwargs.get("prompt", ""))
            calls.append(prompt)
            outline = [
                {"index": 1, "kind": "cover", "title": "Quantum AI Overview", "goals": ["Quantum AI opening"]},
                {"index": 2, "kind": "content", "title": "Quantum AI Market", "goals": ["Quantum AI demand", "Enterprise adoption"]},
                {"index": 3, "kind": "content", "title": "Quantum AI Technology", "goals": ["Hybrid algorithms", "Hardware roadmap"]},
                {"index": 4, "kind": "content", "title": "Quantum AI Risks", "goals": ["Data readiness", "Talent gap"]},
                {"index": 5, "kind": "summary", "title": "Quantum AI Summary", "goals": ["Key takeaways", "Next steps"]},
            ]
            if "outline" in prompt.lower():
                return json.dumps(outline)
            return json.dumps(
                [
                    {"index": row["index"], "kind": row["kind"], "title": row["title"], "bullets": row["goals"], "notes": f"Notes for {row['title']}"}
                    for row in outline
                ]
            )

        with TestClient(app) as client:
            user_id, auth = _register_and_login(client, "ppt_full_path_template_user")
            app.state.repositories.providers.upsert(
                LLMProviderConfig(
                    provider_id="ppt-user-provider",
                    user_id=user_id,
                    provider_type="openai_compatible",
                    display_name="User Model",
                    base_url="https://user-model.example/v1",
                    api_key_encrypted="user-secret",
                    model_name="user-ppt-model",
                    chat_model="user-ppt-model",
                    timeout_seconds=60,
                    is_default=True,
                )
            )
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "ppt",
                    "user_requirement": "Create a 5 page Quantum AI market overview.",
                    "template_choice": str(template_dir),
                    "pages": 5,
                    "style": str(template_dir),
                    "language": "en-US",
                },
                headers=auth,
            )
            assert create.status_code == 200
            create_data = create.json()["data"]
            assert create_data["template_choice"] == "ppt1"
            assert "TemplateChoice=ppt1" in create_data["user_requirement"]
            task_id = create_data["task_id"]

            with patch("app.services.llm_runtime.text_generator.LLMTextGenerator.generate", fake_generate):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False, "require_llm": True}, headers=auth)

            assert run.status_code == 200
            data = run.json()["data"]
            assert data["status"] == "completed"
            assert Path(data["output_path"]).exists()
            assert len(calls) >= 2

            detail = client.get(f"/v1/tasks/{task_id}", headers=auth)
            assert detail.status_code == 200
            detail_data = detail.json()["data"]
            skill_names = [c["skill_name"] for c in detail_data["skill_calls"]]
            assert "llm_text_generation" in skill_names
            assert "ppt_generation_template_render" in skill_names
            render_call = next(c for c in detail_data["skill_calls"] if c["skill_name"] == "ppt_generation_template_render")
            assert '"template_name": "ppt1"' in render_call["input"]
            assert any("template_script_render_succeeded" in e["message"] for e in detail_data["events"])


def test_ppt_generation_uses_templatechoice_even_if_style_is_invalid():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        templates_root = Path(__file__).resolve().parents[1] / "app" / "templates" / "ppt"
        chosen_dir = templates_root / "chosen_template_for_step6"
        try:
            chosen_dir.mkdir(parents=True, exist_ok=True)
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[0])
            prs.save(str(chosen_dir / "template.pptx"))
            (chosen_dir / "render_from_template.py").write_text(
                (
                    "from pathlib import Path\n"
                    "from pptx import Presentation\n\n"
                    "def render(payload, output_path, template_path, meta, rules):\n"
                    "    prs = Presentation(str(template_path))\n"
                    "    while len(prs.slides) > 0:\n"
                    "        slide_id = prs.slides._sldIdLst[0]\n"
                    "        prs.part.drop_rel(slide_id.rId)\n"
                    "        prs.slides._sldIdLst.remove(slide_id)\n"
                    "    for idx, slide_payload in enumerate(payload.get('slides', []), start=1):\n"
                    "        layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None\n"
                    "        slide = prs.slides.add_slide(layout)\n"
                    "        if slide.shapes.title is not None:\n"
                    "            slide.shapes.title.text = str(slide_payload.get('title', f'Slide {idx}'))\n"
                    "    out = Path(output_path)\n"
                    "    out.parent.mkdir(parents=True, exist_ok=True)\n"
                    "    prs.save(str(out))\n"
                    "    return str(out)\n"
                ),
                encoding="utf-8",
            )
            (chosen_dir / "template.meta.json").write_text(
                json.dumps(
                    {
                        "schema_version": "v1",
                        "template_type": "ppt",
                        "template_name": "chosen_template_for_step6",
                        "slide_size": {"width_inches": 13.333, "height_inches": 7.5, "aspect_ratio": "16:9"},
                        "theme": {"name": "default", "palette": {"primary": "#000"}},
                        "layout_map": {"cover": "Title Slide", "content": "Title and Content", "summary": "Title and Content"},
                        "text_style": {"title": {"font": "Arial", "size_pt": 32}, "body": {"font": "Arial", "size_pt": 18}},
                        "text_slots": [
                            {"slide_kind": "cover", "layout_name": "Title Slide", "slot_count": 1, "slots": [{"slot_id": "cover_title", "role": "title", "required": True}]},
                            {"slide_kind": "content", "layout_name": "Title and Content", "slot_count": 1, "slots": [{"slot_id": "content_title", "role": "title", "required": True}]},
                            {"slide_kind": "summary", "layout_name": "Title and Content", "slot_count": 1, "slots": [{"slot_id": "summary_title", "role": "title", "required": True}]},
                        ],
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            (chosen_dir / "template.rules.json").write_text(json.dumps({"schema_version": "v1", "rules": []}), encoding="utf-8")
            (chosen_dir / "template.schema.json").write_text(
                json.dumps({"schema_version": "v1", "template_type": "ppt", "slide_contracts": {"cover": ["cover_title"], "content": ["content_title"], "summary": ["summary_title"]}}),
                encoding="utf-8",
            )

            with TestClient(app) as client:
                user_id, auth = _register_and_login(client, "chosen_templatechoice_user")
                create = client.post(
                    "/v1/tasks",
                    json={
                        "user_id": user_id,
                        "task_type": "ppt",
                        "user_requirement": "Build slides for demo.",
                        "template_choice": "chosen_template_for_step6",
                        "pages": 8,
                        "style": "this_style_does_not_exist",
                        "language": "en-US",
                    },
                    headers=auth,
                )
                assert create.status_code == 200
                create_data = create.json()["data"]
                assert create_data["template_choice"] == "chosen_template_for_step6"
                assert "TemplateChoice=chosen_template_for_step6" in create_data["user_requirement"]
                task_id = create_data["task_id"]
                upload = client.post(
                    f"/v1/tasks/{task_id}/upload",
                    files={"upload": ("sample.txt", BytesIO(b"slide content source text"), "text/plain")},
                    headers=auth,
                )
                assert upload.status_code == 200
                parse = client.post(f"/v1/tasks/{task_id}/parse", json={"force": False}, headers=auth)
                assert parse.status_code == 200
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False}, headers=auth)
                assert run.status_code == 200
                assert run.json()["data"]["status"] == "completed"
                slides_path = Path(run.json()["data"]["slides_path"])
                slides = json.loads(slides_path.read_text(encoding="utf-8"))
                assert slides[0]["layout_intent"]["layout_name"] == "Title Slide"
                assert slides[1]["layout_intent"]["slots"]["body_slot"]["text_style"]["font"] == "Arial"
        finally:
            if chosen_dir.exists():
                shutil.rmtree(chosen_dir)


def test_template_generation_returns_requires_user_completion_for_incomplete_ppt_design():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        templates_root = Path(__file__).resolve().parents[1] / "app" / "templates" / "ppt"
        template_dir = templates_root / "IncompletePptTemplate"
        with TestClient(app) as client:
            user_id, auth = _register_and_login(client, "template_generation_recovery_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "template_generation",
                    "user_requirement": "\n".join(
                        [
                            "Please generate a ppt template from this sample.",
                            "TemplateTarget=ppt",
                            "TemplateName=IncompletePptTemplate",
                            "ForceInvalidBundle=true",
                        ]
                    ),
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
                headers=auth,
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[0])
            stream = BytesIO()
            prs.save(stream)
            stream.seek(0)

            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("template_source.pptx", stream, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                headers=auth,
            )
            assert upload.status_code == 200

            run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False}, headers=auth)
            assert run.status_code == 200
            data = run.json()["data"]
            assert data["status"] == "requires_user_completion"
            assert data["template_type"] == "ppt"
            assert data["template_name"] == "IncompletePptTemplate"
            assert data["missing_items"] == ["template.rules.json"]
            assert data["bundle_validation"]["ok"] is False
            assert isinstance(data.get("missing_fields"), list)
            assert isinstance(data.get("suggested_values"), dict)
            assert data["suggested_values"].get("rules.rules")
            assert data.get("resume_token")
            detail = client.get(f"/v1/tasks/{task_id}", headers=auth)
            assert detail.status_code == 200
            assert detail.json()["data"]["task"]["status"] == "requires_user_completion"
        if template_dir.exists():
            shutil.rmtree(template_dir)


def test_template_generation_failure_updates_task_status_instead_of_stale_running():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, auth = _register_and_login(client, "template_generation_failure_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "template_generation",
                    "user_requirement": "\n".join(
                        [
                            "Please generate a ppt template from this sample.",
                            "TemplateTarget=ppt",
                            "TemplateName=ExplodingTemplate",
                        ]
                    ),
                },
                headers=auth,
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[0])
            stream = BytesIO()
            prs.save(stream)
            stream.seek(0)
            upload = client.post(
                f"/v1/tasks/{task_id}/upload",
                files={"upload": ("template_source.pptx", stream, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                headers=auth,
            )
            assert upload.status_code == 200

            with patch("app.agents.task_agents.template_generation_task_agent.TemplateGenerationTaskAgent.execute", side_effect=RuntimeError("boom")):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False}, headers=auth)

            assert run.status_code == 400
            detail = client.get(f"/v1/tasks/{task_id}", headers=auth)
            assert detail.status_code == 200
            detail_data = detail.json()["data"]
            assert detail_data["task"]["status"] == "failed_generation"
            assert any("template_generation_failed=boom" in e["message"] for e in detail_data["events"])
            assert not any("stale_running_task_auto_closed" in e["message"] for e in detail_data["events"])


def test_infer_task_type_returns_generic_when_no_keyword():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            _register_and_login(client, "task_flow_infer_type_user")
            resp = client.post(
                "/v1/tasks/infer-type",
                json={"requirement": "Please help me process this request with custom workflow."},
            )
            assert resp.status_code == 200
            assert resp.json()["data"]["task_type"] == "generic_task"


def test_auto_task_type_routes_to_generic_task_and_runs():
    with TemporaryDirectory() as temp_dir:
        settings.data_dir = Path(temp_dir)
        app = create_app()
        with TestClient(app) as client:
            user_id, _ = _register_and_login(client, "task_flow_auto_user")
            create = client.post(
                "/v1/tasks",
                json={
                    "user_id": user_id,
                    "task_type": "auto",
                    "user_requirement": "Please handle this custom internal workflow without special format.",
                    "pages": 8,
                    "style": "academic_simple",
                    "language": "en-US",
                },
            )
            assert create.status_code == 200
            task_id = create.json()["data"]["task_id"]

            detail = client.get(f"/v1/tasks/{task_id}")
            assert detail.status_code == 200
            assert detail.json()["data"]["task"]["task_type"] == "generic_task"

            with patch(
                "app.services.llm_runtime.text_generator.LLMTextGenerator.generate",
                return_value="# Generic Task Output\n\n## Summary\nDone.",
            ):
                run = client.post(f"/v1/tasks/{task_id}/run", json={"rerun": False})
                assert run.status_code == 200
                assert run.json()["data"]["status"] == "completed"
                assert run.json()["data"]["task_type"] == "generic_task"
                output_path = Path(run.json()["data"]["output_path"])
                assert output_path.exists()
                assert output_path.suffix == ".md"
