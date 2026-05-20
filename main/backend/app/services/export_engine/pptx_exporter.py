import json
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


class PptxExportError(Exception):
    pass


class PptxExporter:
    def export(
        self,
        slides: list[dict[str, Any]],
        output_path: Path,
        template_path: Optional[Path] = None,
        template_bundle: Optional[dict[str, Any]] = None,
    ) -> Path:
        try:
            normalized_slides = [self._normalize_slide_payload(slide, index=i) for i, slide in enumerate(slides, start=1)]
            effective_template_path = template_path
            if template_bundle and template_bundle.get("template_file"):
                effective_template_path = Path(str(template_bundle["template_file"]))
            meta, rules = self._load_bundle_meta_rules(template_bundle)
            rule_settings = self._resolve_rule_settings(rules)
            style_mode = self._resolve_style_mode(meta=meta, rules=rules)
            background_policy = self._resolve_background_policy(meta=meta, rules=rules)
            palette = self._resolve_palette(meta)
            title_style = self._resolve_text_style(meta, "title", default_size=30.0, default_color=palette.get("primary") or palette.get("text"))
            body_style = self._resolve_text_style(meta, "body", default_size=18.0, default_color=palette.get("text"))

            if effective_template_path is not None and effective_template_path.exists():
                prs = Presentation(str(effective_template_path))
                self._clear_existing_slides(prs)
            else:
                prs = Presentation()
            self._apply_slide_size_from_meta(prs, meta)
            title_layout = self._pick_layout(prs, (meta.get("layout_map") or {}).get("cover"), fallback=0)
            content_layout = self._pick_layout(prs, (meta.get("layout_map") or {}).get("content"), fallback=1 if len(prs.slide_layouts) > 1 else 0)
            blank_layout = self._pick_layout(prs, (meta.get("layout_map") or {}).get("blank"), fallback=6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1)

            for index, slide in enumerate(normalized_slides):
                kind = slide.get("kind", "content")
                image_placeholders = slide.get("image_placeholders", []) or []
                if index == 0 or kind == "cover":
                    s = prs.slides.add_slide(title_layout)
                    self._apply_slide_background(s, palette, background_policy=background_policy)
                    self._fill_title_slot(
                        s,
                        str(slide.get("title", "Cover")),
                        title_style=title_style,
                        rule_settings=rule_settings,
                        style_mode=style_mode,
                    )
                    self._fill_subtitle_slot(
                        s,
                        [str(x) for x in (slide.get("bullets", [])[:2])],
                        body_style=body_style,
                        rule_settings=rule_settings,
                        style_mode=style_mode,
                    )
                elif image_placeholders:
                    s = prs.slides.add_slide(blank_layout)
                    self._apply_slide_background(s, palette, background_policy=background_policy)
                    self._fill_title_slot(
                        s,
                        str(slide.get("title", f"Slide {index + 1}")),
                        title_style=title_style,
                        rule_settings=rule_settings,
                        style_mode=style_mode,
                    )
                    self._fill_body_slot(
                        s,
                        [str(x) for x in (slide.get("bullets", []) or ["TBD"])],
                        body_style=body_style,
                        rule_settings=rule_settings,
                        allow_fallback=True,
                        style_mode=style_mode,
                    )
                    self._fill_image_slot(
                        s,
                        image_placeholders=image_placeholders,
                        body_style=body_style,
                        rule_settings=rule_settings,
                        palette=palette,
                        style_mode=style_mode,
                    )
                else:
                    s = prs.slides.add_slide(content_layout)
                    self._apply_slide_background(s, palette, background_policy=background_policy)
                    self._fill_title_slot(
                        s,
                        str(slide.get("title", f"Slide {index + 1}")),
                        title_style=title_style,
                        rule_settings=rule_settings,
                        style_mode=style_mode,
                    )
                    self._fill_body_slot(
                        s,
                        [str(x) for x in (slide.get("bullets", []) or ["TBD"])],
                        body_style=body_style,
                        rule_settings=rule_settings,
                        allow_fallback=True,
                        style_mode=style_mode,
                    )

                notes_text = self._prepare_text(str(slide.get("notes", "")), role="notes", rule_settings=rule_settings)
                s.notes_slide.notes_text_frame.text = notes_text

            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
            return output_path
        except Exception as exc:
            raise PptxExportError(f"PPTX export failed: {exc}") from exc

    def _normalize_slide_payload(self, slide: dict[str, Any], *, index: int) -> dict[str, Any]:
        kind = str(slide.get("kind", "content")).strip() or "content"
        title = (
            str(slide.get("title", "")).strip()
            or str(slide.get("header", "")).strip()
            or str(slide.get("topic", "")).strip()
            or f"Slide {index}"
        )
        bullets_raw = slide.get("bullets", None)
        if not isinstance(bullets_raw, list):
            for fallback_key in ("content", "body", "points"):
                candidate = slide.get(fallback_key)
                if isinstance(candidate, list):
                    bullets_raw = candidate
                    break
                if isinstance(candidate, str) and candidate.strip():
                    bullets_raw = [row.strip("- ").strip() for row in candidate.splitlines() if row.strip()]
                    break
        bullets = [str(item).strip() for item in (bullets_raw or []) if str(item).strip()]
        if not bullets:
            bullets = ["TBD"]
        notes = str(slide.get("notes", "")).strip() or str(slide.get("speaker_notes", "")).strip()
        image_placeholders = slide.get("image_placeholders", [])
        if not isinstance(image_placeholders, list):
            image_placeholders = []
        return {
            "index": int(slide.get("index", index) or index),
            "kind": kind,
            "title": title,
            "bullets": bullets[:8],
            "notes": notes,
            "image_placeholders": image_placeholders,
        }

    def _pick_layout(self, prs: Presentation, preferred_name: Any, fallback: int):
        name = str(preferred_name or "").strip().lower()
        if name:
            for layout in prs.slide_layouts:
                if str(getattr(layout, "name", "") or "").strip().lower() == name:
                    return layout
        idx = max(0, min(fallback, len(prs.slide_layouts) - 1))
        return prs.slide_layouts[idx]

    def _load_bundle_meta_rules(self, template_bundle: Optional[dict[str, Any]]) -> tuple[dict[str, Any], dict[str, Any]]:
        if not template_bundle:
            return {}, {}
        meta_path = Path(str(template_bundle.get("meta_path", ""))) if template_bundle.get("meta_path") else None
        rules_path = Path(str(template_bundle.get("rules_path", ""))) if template_bundle.get("rules_path") else None
        meta: dict[str, Any] = {}
        rules: dict[str, Any] = {}
        if meta_path and meta_path.exists():
            try:
                loaded = json.loads(meta_path.read_text(encoding="utf-8"))
                if isinstance(loaded, dict):
                    meta = loaded
            except Exception:
                meta = {}
        if rules_path and rules_path.exists():
            try:
                loaded = json.loads(rules_path.read_text(encoding="utf-8"))
                if isinstance(loaded, dict):
                    rules = loaded
            except Exception:
                rules = {}
        return meta, rules

    def _resolve_palette(self, meta: dict[str, Any]) -> dict[str, str]:
        theme = meta.get("theme") if isinstance(meta, dict) else {}
        palette = theme.get("palette") if isinstance(theme, dict) else {}
        return {str(k): str(v) for k, v in palette.items()} if isinstance(palette, dict) else {}

    def _resolve_style_mode(self, *, meta: dict[str, Any], rules: dict[str, Any]) -> str:
        render = meta.get("render") if isinstance(meta, dict) else {}
        mode = str((render or {}).get("style_mode") or rules.get("style_mode") or "inherit_first").strip().lower()
        return mode if mode in {"inherit_first", "force_meta"} else "inherit_first"

    def _resolve_background_policy(self, *, meta: dict[str, Any], rules: dict[str, Any]) -> dict[str, Any]:
        render_meta = meta.get("render") if isinstance(meta, dict) else {}
        render_rules = rules.get("render") if isinstance(rules, dict) else {}
        merged: dict[str, Any] = {}
        if isinstance(render_meta, dict):
            merged.update(render_meta)
        if isinstance(render_rules, dict):
            merged.update(render_rules)
        preserve = bool(merged.get("preserve_background", True))
        force_background_color = str(merged.get("force_background_color") or "").strip()
        apply_palette_when_missing = bool(merged.get("apply_palette_background_when_missing", False))
        return {
            "preserve_background": preserve,
            "force_background_color": force_background_color,
            "apply_palette_background_when_missing": apply_palette_when_missing,
        }

    def _resolve_text_style(self, meta: dict[str, Any], role: str, *, default_size: float, default_color: Optional[str]) -> dict[str, Any]:
        text_style = meta.get("text_style") if isinstance(meta, dict) else {}
        raw = text_style.get(role) if isinstance(text_style, dict) else {}
        raw = raw if isinstance(raw, dict) else {}
        try:
            size = float(raw.get("size_pt")) if raw.get("size_pt") is not None else default_size
        except Exception:
            size = default_size
        try:
            line_spacing = float(raw.get("line_spacing")) if raw.get("line_spacing") is not None else None
        except Exception:
            line_spacing = None
        try:
            indent = float(raw.get("indent")) if raw.get("indent") is not None else 0.0
        except Exception:
            indent = 0.0
        return {
            "font": str(raw.get("font")) if raw.get("font") else None,
            "size_pt": size if raw.get("size_pt") is not None else None,
            "fallback_size_pt": default_size,
            "line_spacing": line_spacing,
            "indent": indent,
            "color": str(raw.get("color")) if raw.get("color") else None,
            "fallback_color": str(default_color or ""),
        }

    def _resolve_rule_settings(self, rules: dict[str, Any]) -> dict[str, Any]:
        settings: dict[str, Any] = {
            "overflow_action": "truncate",
            "title_max_chars": 120,
            "body_max_chars": 220,
            "notes_max_chars": 1200,
            "role_actions": {},
        }
        rows = rules.get("rules", []) if isinstance(rules, dict) else []
        if not isinstance(rows, list):
            return settings
        for row in rows:
            if not isinstance(row, dict):
                continue
            name = str(row.get("name", "")).strip().lower()
            if name == "text_overflow":
                action = str(row.get("action", "")).strip().lower()
                if action in {"truncate", "shrink"}:
                    settings["overflow_action"] = action
                continue
            if name in {"title_max_chars", "body_max_chars", "notes_max_chars"}:
                try:
                    settings[name] = max(4, int(row.get("max_chars")))
                except Exception:
                    pass
                action = str(row.get("action", "")).strip().lower()
                if action in {"truncate", "shrink"}:
                    role = name.split("_", 1)[0]
                    settings["role_actions"][role] = action
        return settings

    def _prepare_text(self, text: str, *, role: str, rule_settings: dict[str, Any]) -> str:
        max_len = int(rule_settings.get(f"{role}_max_chars") or rule_settings.get("body_max_chars") or 220)
        action = (rule_settings.get("role_actions") or {}).get(role) or rule_settings.get("overflow_action") or "truncate"
        if len(text) <= max_len:
            return text
        if action == "truncate":
            return text[: max_len - 3].rstrip() + "..."
        return text

    def _font_size_for_text(self, base_size: float, source_text: str, *, role: str, rule_settings: dict[str, Any]) -> float:
        max_len = int(rule_settings.get(f"{role}_max_chars") or rule_settings.get("body_max_chars") or 220)
        action = (rule_settings.get("role_actions") or {}).get(role) or rule_settings.get("overflow_action") or "truncate"
        if action != "shrink" or len(source_text) <= max_len:
            return base_size
        ratio = max_len / max(1, len(source_text))
        return max(10.0, base_size * max(0.55, ratio))

    def _apply_shape_text_style(
        self,
        shape: Any,
        style: dict[str, Any],
        *,
        source_text: str,
        rule_settings: dict[str, Any],
        role: str,
        style_mode: str,
    ) -> None:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            return
        for para in shape.text_frame.paragraphs:
            self._apply_paragraph_text_style(para, style, source_text=source_text, rule_settings=rule_settings, role=role, style_mode=style_mode)

    def _apply_paragraph_text_style(
        self,
        para: Any,
        style: dict[str, Any],
        *,
        source_text: str,
        rule_settings: dict[str, Any],
        role: str,
        style_mode: str,
    ) -> None:
        para_style_font = style.get("font")
        para_style_size = style.get("size_pt")
        para_style_color = style.get("color")
        inherit_first = style_mode == "inherit_first"
        if inherit_first:
            if style.get("line_spacing") is not None:
                para.line_spacing = float(style["line_spacing"])
            indent = float(style.get("indent") or 0.0)
            if indent:
                para.left_margin = Pt(indent)
            return

        current_name = getattr(para.font, "name", None)
        current_size = getattr(para.font, "size", None)
        if (not inherit_first) or current_name is None:
            if para_style_font:
                para.font.name = str(para_style_font)
        if (not inherit_first) or current_size is None:
            if para_style_size is not None:
                base = float(para_style_size)
            else:
                base = float(style.get("fallback_size_pt") or 18.0)
            para.font.size = Pt(self._font_size_for_text(base, source_text, role=role, rule_settings=rule_settings))

        current_color = getattr(getattr(para.font, "color", None), "rgb", None)
        if (not inherit_first) or current_color is None:
            color_text = str(para_style_color or style.get("fallback_color") or "")
            color = self._parse_hex_color(color_text)
            if color is not None:
                para.font.color.rgb = color
        if style.get("line_spacing") is not None:
            para.line_spacing = float(style["line_spacing"])
        indent = float(style.get("indent") or 0.0)
        if indent:
            para.left_margin = Pt(indent)

    def _parse_hex_color(self, value: str) -> Optional[RGBColor]:
        text = (value or "").strip().lstrip("#")
        if len(text) != 6:
            return None
        try:
            return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
        except Exception:
            return None

    def _apply_slide_size_from_meta(self, prs: Presentation, meta: dict[str, Any]) -> None:
        slide_size = meta.get("slide_size") if isinstance(meta, dict) else {}
        if not isinstance(slide_size, dict):
            return
        try:
            width = float(slide_size.get("width_inches"))
            height = float(slide_size.get("height_inches"))
        except Exception:
            return
        if width > 0 and height > 0:
            prs.slide_width = Inches(width)
            prs.slide_height = Inches(height)

    def _apply_slide_background(self, slide_obj: Any, palette: dict[str, str], *, background_policy: dict[str, Any]) -> None:
        if bool(background_policy.get("preserve_background", True)):
            return
        force_color = str(background_policy.get("force_background_color") or "").strip()
        if force_color:
            background = self._parse_hex_color(force_color)
        elif bool(background_policy.get("apply_palette_background_when_missing", False)):
            background = self._parse_hex_color(palette.get("background", ""))
        else:
            background = None
        if background is None:
            return
        fill = slide_obj.background.fill
        fill.solid()
        fill.fore_color.rgb = background

    def _clear_existing_slides(self, prs: Presentation) -> None:
        slide_id_list = prs.slides._sldIdLst
        slides = list(slide_id_list)
        for slide_id in slides:
            rel_id = slide_id.rId
            prs.part.drop_rel(rel_id)
            slide_id_list.remove(slide_id)

    def _find_slot(self, slide_obj: Any, role: str) -> Optional[Any]:
        role_key = (role or "").strip().lower()
        title_shape = getattr(slide_obj.shapes, "title", None)
        if role_key == "title" and title_shape is not None and hasattr(title_shape, "text_frame"):
            return title_shape
        candidates: list[Any] = []
        for ph in getattr(slide_obj, "placeholders", []):
            if not hasattr(ph, "placeholder_format"):
                continue
            pf = ph.placeholder_format
            ptype = str(getattr(pf, "type", "")).lower()
            name = str(getattr(ph, "name", "")).lower()
            if role_key == "subtitle":
                if "subtitle" in ptype or "subtitle" in name:
                    return ph
                continue
            if role_key == "body":
                if not hasattr(ph, "text_frame"):
                    continue
                if "body" in ptype or "text" in ptype or "object" in ptype:
                    candidates.append(ph)
                    continue
                if any(token in name for token in ("content", "body", "text")):
                    candidates.append(ph)
                    continue
            if role_key == "image":
                if "pic" in ptype or "media" in ptype or "picture" in ptype or "image" in name:
                    return ph
        if role_key == "body" and candidates:
            candidates.sort(key=lambda shape: int(getattr(shape.placeholder_format, "idx", 9999)))
            return candidates[0]
        return None

    def _fill_title_slot(
        self,
        slide_obj: Any,
        raw_title: str,
        *,
        title_style: dict[str, Any],
        rule_settings: dict[str, Any],
        style_mode: str,
    ) -> None:
        slot = self._find_slot(slide_obj, "title")
        if slot is None or not hasattr(slot, "text_frame"):
            return
        title_text = self._prepare_text(str(raw_title), role="title", rule_settings=rule_settings)
        slot.text = title_text
        self._apply_shape_text_style(slot, title_style, source_text=title_text, rule_settings=rule_settings, role="title", style_mode=style_mode)

    def _fill_subtitle_slot(
        self,
        slide_obj: Any,
        rows: list[str],
        *,
        body_style: dict[str, Any],
        rule_settings: dict[str, Any],
        style_mode: str,
    ) -> None:
        slot = self._find_slot(slide_obj, "subtitle")
        if slot is None or not hasattr(slot, "text_frame"):
            return
        subtitle_text = "\n".join([self._prepare_text(str(row), role="body", rule_settings=rule_settings) for row in rows if str(row).strip()])
        if not subtitle_text:
            return
        slot.text = subtitle_text
        self._apply_shape_text_style(slot, body_style, source_text=subtitle_text, rule_settings=rule_settings, role="body", style_mode=style_mode)

    def _fill_body_slot(
        self,
        slide_obj: Any,
        bullets: list[str],
        *,
        body_style: dict[str, Any],
        rule_settings: dict[str, Any],
        allow_fallback: bool,
        style_mode: str,
    ) -> bool:
        slot = self._find_slot(slide_obj, "body")
        if slot is not None and hasattr(slot, "text_frame"):
            tf = slot.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                text = self._prepare_text(str(bullet), role="body", rule_settings=rule_settings)
                p.text = text
                self._apply_paragraph_text_style(p, body_style, source_text=text, rule_settings=rule_settings, role="body", style_mode=style_mode)
            return True
        if not allow_fallback:
            return False
        body_box = slide_obj.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.8))
        tf = body_box.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            text = self._prepare_text(str(bullet), role="body", rule_settings=rule_settings)
            p.text = text
            p.level = 0
            self._apply_paragraph_text_style(p, body_style, source_text=text, rule_settings=rule_settings, role="body", style_mode="force_meta")
        return False

    def _fill_image_slot(
        self,
        slide_obj: Any,
        *,
        image_placeholders: list[dict[str, Any]],
        body_style: dict[str, Any],
        rule_settings: dict[str, Any],
        palette: dict[str, str],
        style_mode: str,
    ) -> None:
        slot = self._find_slot(slide_obj, "image")
        has_template_image_slot = slot is not None
        ph = (image_placeholders or [{}])[0]
        label = str(ph.get("label", "Image placeholder"))[:120]
        source = str(ph.get("source", "unknown"))[:200]
        lines = [
            "IMAGE PLACEHOLDER",
            f"Label: {label}",
            f"Source: {source}",
        ]
        if slot is not None and hasattr(slot, "text_frame"):
            tf = slot.text_frame
            tf.clear()
            for i, row in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = row
                self._apply_paragraph_text_style(p, body_style, source_text=row, rule_settings=rule_settings, role="body", style_mode=style_mode)
            return
        if has_template_image_slot:
            # Template already defines an image slot (often picture placeholder without text frame).
            # Keep shape/effects unchanged and write the placeholder note to an existing body text slot.
            if self._append_image_placeholder_note_to_body_slot(
                slide_obj,
                lines=lines,
                body_style=body_style,
                rule_settings=rule_settings,
                style_mode=style_mode,
            ):
                return
        self._render_slide_with_image_placeholder(
            slide_obj,
            {"image_placeholders": image_placeholders},
            title_style={},
            body_style=body_style,
            rule_settings=rule_settings,
            palette=palette,
            skip_title_and_body=True,
            style_mode=style_mode,
        )

    def _append_image_placeholder_note_to_body_slot(
        self,
        slide_obj: Any,
        *,
        lines: list[str],
        body_style: dict[str, Any],
        rule_settings: dict[str, Any],
        style_mode: str,
    ) -> bool:
        body_slot = self._find_slot(slide_obj, "body")
        if body_slot is None or not hasattr(body_slot, "text_frame"):
            return False
        tf = body_slot.text_frame
        for row in lines:
            if not str(row).strip():
                continue
            p = tf.add_paragraph()
            p.text = row
            self._apply_paragraph_text_style(
                p,
                body_style,
                source_text=row,
                rule_settings=rule_settings,
                role="body",
                style_mode=style_mode,
            )
        return True

    def _render_slide_with_image_placeholder(
        self,
        slide_obj,
        slide_data: dict[str, Any],
        *,
        title_style: dict[str, Any],
        body_style: dict[str, Any],
        rule_settings: dict[str, Any],
        palette: dict[str, str],
        skip_title_and_body: bool = False,
        style_mode: str = "force_meta",
    ) -> None:
        if not skip_title_and_body:
            title_box = slide_obj.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.8))
            tf_title = title_box.text_frame
            title_text = self._prepare_text(str(slide_data.get("title", "Content")), role="title", rule_settings=rule_settings)
            tf_title.text = title_text
            self._apply_paragraph_text_style(
                tf_title.paragraphs[0],
                title_style,
                source_text=title_text,
                rule_settings=rule_settings,
                role="title",
                style_mode=style_mode,
            )

            left_box = slide_obj.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(6.1), Inches(5.2))
            tf_left = left_box.text_frame
            tf_left.clear()
            bullets = slide_data.get("bullets", []) or ["TBD"]
            for i, bullet in enumerate(bullets[:6]):
                p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
                text = self._prepare_text(str(bullet), role="body", rule_settings=rule_settings)
                p.text = text
                p.level = 0
                self._apply_paragraph_text_style(p, body_style, source_text=text, rule_settings=rule_settings, role="body", style_mode=style_mode)

        # Right image reserved area
        image_shape = slide_obj.shapes.add_shape(
            autoshape_type_id=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left=Inches(7.1),
            top=Inches(1.6),
            width=Inches(5.6),
            height=Inches(4.7),
        )
        image_shape.fill.background()
        image_shape.line.width = Pt(1.5)
        accent = self._parse_hex_color(palette.get("accent", "") or palette.get("secondary", ""))
        if accent is not None:
            image_shape.line.color.rgb = accent

        ph = (slide_data.get("image_placeholders") or [{}])[0]
        label = str(ph.get("label", "Image placeholder"))[:120]
        source = str(ph.get("source", "unknown"))[:200]
        info_box = slide_obj.shapes.add_textbox(Inches(7.3), Inches(1.9), Inches(5.2), Inches(4.2))
        tf_info = info_box.text_frame
        tf_info.clear()
        p0 = tf_info.paragraphs[0]
        p0.text = "IMAGE PLACEHOLDER"
        self._apply_paragraph_text_style(p0, body_style, source_text=p0.text, rule_settings=rule_settings, role="body", style_mode=style_mode)
        p1 = tf_info.add_paragraph()
        p1.text = f"Label: {label}"
        self._apply_paragraph_text_style(p1, body_style, source_text=p1.text, rule_settings=rule_settings, role="body", style_mode=style_mode)
        p2 = tf_info.add_paragraph()
        p2.text = f"Source: {source}"
        self._apply_paragraph_text_style(p2, body_style, source_text=p2.text, rule_settings=rule_settings, role="body", style_mode=style_mode)
