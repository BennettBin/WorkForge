from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Callable


class ParseError(Exception):
    pass


@dataclass
class ParseResult:
    text: str
    used_ocr: bool = False
    parser_name: str = ""


def _parse_txt(path: Path) -> ParseResult:
    text = path.read_text(encoding="utf-8", errors="ignore").strip()
    if not text:
        raise ParseError("TXT file contains no readable text.")
    return ParseResult(text=text, parser_name="txt")


def _parse_docx(path: Path) -> ParseResult:
    try:
        from docx import Document

        doc = Document(str(path))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        text = "\n".join(paragraphs).strip()
        if not text:
            raise ParseError("DOCX file contains no readable text.")
        return ParseResult(text=text, parser_name="docx")
    except ParseError:
        raise
    except Exception as exc:
        raise ParseError("DOCX file parsing failed (file may be damaged).") from exc


def _parse_pptx(path: Path) -> ParseResult:
    from pptx import Presentation

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise ParseError("PPT parsing failed. Legacy .ppt may require conversion to .pptx.") from exc
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text.strip())
    text = "\n".join([c for c in chunks if c]).strip()
    if not text:
        raise ParseError("PPT/PPTX file contains no readable text.")
    return ParseResult(text=text, parser_name="pptx")


def _try_pdf_ocr(path: Path) -> ParseResult:
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except Exception as exc:
        raise ParseError("Scanned PDF detected, OCR dependencies are unavailable.") from exc

    images = convert_from_path(str(path))
    chunks = [pytesseract.image_to_string(img) for img in images]
    text = "\n".join([c.strip() for c in chunks if c and c.strip()]).strip()
    if not text:
        raise ParseError("OCR completed but no text was extracted from scanned PDF.")
    return ParseResult(text=text, used_ocr=True, parser_name="pdf_ocr")


def _parse_pdf(path: Path) -> ParseResult:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        chunks: list[str] = []
        for page in reader.pages:
            chunks.append((page.extract_text() or "").strip())
        text = "\n".join([c for c in chunks if c]).strip()
        if text:
            return ParseResult(text=text, parser_name="pdf")
        return _try_pdf_ocr(path)
    except ParseError:
        raise
    except Exception as exc:
        raise ParseError("PDF parsing failed (file may be damaged or encrypted).") from exc


def _parse_doc(path: Path) -> ParseResult:
    # MVP fallback for legacy .doc: best-effort decode with explicit warning path.
    raw = path.read_bytes()
    text = raw.decode("utf-8", errors="ignore").strip()
    if not text:
        text = raw.decode("latin-1", errors="ignore").strip()
    if not text:
        raise ParseError("DOC file could not be parsed. Please convert to DOCX.")
    return ParseResult(text=text, parser_name="doc_fallback")


def _parse_xlsx(path: Path) -> ParseResult:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), data_only=True, read_only=True)
        try:
            ws = wb.active
            rows = list(ws.iter_rows(values_only=True))
        finally:
            wb.close()
        if not rows:
            raise ParseError("Excel file is empty.")
        header = [str(c).strip() if c is not None else "" for c in rows[0]]
        lines: list[str] = []
        lines.append("Excel sheet parsed")
        lines.append("Columns: " + ", ".join([h for h in header if h]))
        sample_rows = rows[1:21]
        for idx, row in enumerate(sample_rows, start=1):
            pairs = []
            for c_i, cell in enumerate(row):
                key = header[c_i] if c_i < len(header) and header[c_i] else f"col_{c_i+1}"
                val = "" if cell is None else str(cell).strip()
                if val:
                    pairs.append(f"{key}={val}")
            if pairs:
                lines.append(f"row_{idx}: " + "; ".join(pairs))
        text = "\n".join(lines).strip()
        if not text:
            raise ParseError("Excel file contains no readable rows.")
        return ParseResult(text=text, parser_name="xlsx")
    except ParseError:
        raise
    except Exception as exc:
        raise ParseError("Excel file parsing failed (file may be damaged).") from exc


_PARSERS: dict[str, Callable[[Path], ParseResult]] = {
    "txt": _parse_txt,
    "docx": _parse_docx,
    "doc": _parse_doc,
    "pptx": _parse_pptx,
    "ppt": _parse_pptx,
    "pdf": _parse_pdf,
    "xlsx": _parse_xlsx,
    "xls": _parse_xlsx,
}


def parse_file(path: Path, file_type: str) -> ParseResult:
    parser = _PARSERS.get(file_type.lower())
    if parser is None:
        raise ParseError(f"Unsupported file type for parser: {file_type}")
    return parser(path)
